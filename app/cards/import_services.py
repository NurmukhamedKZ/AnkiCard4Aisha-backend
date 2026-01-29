"""
Import services for various file formats (CSV, Anki, PowerPoint)
and text-based card generation.
"""
import logging
import io
import zipfile
import sqlite3
import asyncio
from typing import List, Dict, Optional
from concurrent.futures import ThreadPoolExecutor

from google import genai
from google.genai import types

from app.config import get_settings

logger = logging.getLogger(__name__)
settings = get_settings()

# Configure Gemini with HTTP timeout
client = genai.Client(
    api_key=settings.GEMINI_API_KEY,
    http_options={"timeout": 120}  # 2 minute timeout for API calls
)

# Thread pool for running sync Gemini calls
_executor = ThreadPoolExecutor(max_workers=4)


def parse_csv_cards(cards_data: List[Dict[str, str]]) -> List[Dict[str, str]]:
    """
    Parse pre-parsed card data from frontend.
    Frontend sends: [{"front": "Q1", "back": "A1"}, ...]
    Returns: [{"question": "Q1", "answer": "A1"}, ...]
    """
    parsed_cards = []
    for card in cards_data:
        if "front" in card and "back" in card:
            question = card["front"].strip()
            answer = card["back"].strip()
            if question and answer:
                parsed_cards.append({
                    "question": question,
                    "answer": answer
                })
    return parsed_cards


TEXT_GENERATION_PROMPT = """
You are a flashcard generator. Given text content, create high-quality question-answer flashcards.

Rules:
- Extract key facts, concepts, and important information
- Create clear, specific questions
- Provide concise, accurate answers
- Each card should test one concept
- Format: question:answer (separated by colon)
- One card per line

Example format:
What is the capital of France?:Paris
Who wrote Romeo and Juliet?:William Shakespeare
What is photosynthesis?:The process by which plants convert light energy into chemical energy
"""


def _generate_cards_from_text_sync(text: str) -> Optional[str]:
    """
    Synchronous Gemini call for text generation (runs in thread pool).
    """
    try:
        logger.info("Generating cards from text using Gemini...")
        response = client.models.generate_content(
            model="gemini-2.0-flash",
            config=types.GenerateContentConfig(system_instruction=TEXT_GENERATION_PROMPT),
            contents=text
        )
        
        if not response or not response.text:
            logger.warning("No response from Gemini")
            return None
        
        return response.text
    except Exception as e:
        logger.error(f"Error in Gemini call: {e}")
        return None


async def generate_cards_from_text(text: str, timeout: float = 180.0) -> List[Dict[str, str]]:
    """
    Generate flashcards from plain text using Gemini AI.
    
    Args:
        text: Plain text content to generate cards from
        timeout: Maximum time to wait for Gemini response (default 180 seconds)
        
    Returns:
        List of card dictionaries with 'question' and 'answer' keys
    """
    if not text.strip():
        return []
    
    try:
        loop = asyncio.get_event_loop()
        response_text = await asyncio.wait_for(
            loop.run_in_executor(_executor, _generate_cards_from_text_sync, text),
            timeout=timeout
        )
        
        if not response_text:
            return []
        
        # Remove markdown code blocks if present
        if response_text.startswith("```"):
            lines = response_text.split("\n")
            response_text = "\n".join(lines[1:-1])
        
        # Parse response into cards
        cards = []
        for line in response_text.split("\n"):
            if ":" in line:
                parts = line.split(":", 1)
                if len(parts) == 2:
                    question = parts[0].strip()
                    answer = parts[1].strip()
                    if question and answer:
                        cards.append({
                            "question": question,
                            "answer": answer
                        })
        
        logger.info(f"Generated {len(cards)} cards from text")
        return cards
        
    except asyncio.TimeoutError:
        logger.error(f"Gemini text generation timed out after {timeout} seconds")
        raise Exception("AI generation timed out. Please try again with less content.")
    except Exception as e:
        logger.error(f"Error generating cards from text: {e}")
        raise


async def extract_and_generate_from_pptx(pptx_bytes: bytes) -> List[Dict[str, str]]:
    """
    Extract text from PowerPoint file and generate flashcards.
    
    Args:
        pptx_bytes: PowerPoint file content
        
    Returns:
        List of card dictionaries
    """
    try:
        from pptx import Presentation
        
        # Load presentation from bytes
        pptx_file = io.BytesIO(pptx_bytes)
        prs = Presentation(pptx_file)
        
        # Extract text from all slides
        text_content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            if slide_text:
                text_content.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
        
        if not text_content:
            logger.warning("No text found in PowerPoint file")
            return []
        
        full_text = "\n\n".join(text_content)
        logger.info(f"Extracted text from {len(prs.slides)} slides")
        
        # Generate cards from extracted text
        return await generate_cards_from_text(full_text)
        
    except ImportError:
        logger.error("python-pptx not installed")
        raise Exception("PowerPoint support not available. Please install python-pptx.")
    except Exception as e:
        logger.error(f"Error processing PowerPoint file: {e}")
        raise


def parse_anki_package(apkg_bytes: bytes) -> Dict[str, List[Dict[str, str]]]:
    """
    Parse Anki .apkg file and extract decks with cards.
    
    Args:
        apkg_bytes: .apkg file content
        
    Returns:
        Dictionary mapping deck names to lists of cards
    """
    try:
        # .apkg is a zip file containing collection.anki2 (SQLite database)
        apkg_file = io.BytesIO(apkg_bytes)
        
        with zipfile.ZipFile(apkg_file, 'r') as zip_ref:
            # Extract collection.anki2
            if 'collection.anki2' not in zip_ref.namelist():
                raise ValueError("Invalid Anki package: collection.anki2 not found")
            
            db_bytes = zip_ref.read('collection.anki2')
            
        # Write to temporary file for SQLite
        db_file = io.BytesIO(db_bytes)
        
        # Connect to SQLite database
        conn = sqlite3.connect(':memory:')
        conn.executescript(db_bytes.decode('utf-8', errors='ignore'))
        cursor = conn.cursor()
        
        # Query decks and cards
        # Anki database schema: cards -> notes -> decks
        cursor.execute("""
            SELECT n.flds, c.did
            FROM cards c
            JOIN notes n ON c.nid = n.id
        """)
        
        rows = cursor.fetchall()
        
        # Group cards by deck
        decks = {}
        for flds, did in rows:
            # Fields are separated by \x1f
            fields = flds.split('\x1f')
            if len(fields) >= 2:
                question = fields[0].strip()
                answer = fields[1].strip()
                
                deck_name = f"Deck_{did}"  # Simplified deck naming
                if deck_name not in decks:
                    decks[deck_name] = []
                
                decks[deck_name].append({
                    "question": question,
                    "answer": answer
                })
        
        conn.close()
        
        logger.info(f"Parsed {len(decks)} decks from Anki package")
        return decks
        
    except Exception as e:
        logger.error(f"Error parsing Anki package: {e}")
        raise
